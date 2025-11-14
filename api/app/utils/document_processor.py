"""
Processador de documentos
Suporte: PDF, CSV, Excel, TXT, Imagens (OCR), Word, PowerPoint
"""
import io
import logging
import subprocess
import tempfile
import os
from typing import Optional
from pathlib import Path

# Básico
import PyPDF2
import pdfplumber
import pandas as pd

# OCR
import pytesseract
from PIL import Image

# Word e PowerPoint
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Extrai texto de diferentes formatos de documento"""
    
    def extract_text(self, content: bytes, filename: str) -> str:
        """Extrai texto do documento baseado na extensão"""
        extension = filename.lower().split('.')[-1]
        
        try:
            # PDFs
            if extension == 'pdf':
                return self._extract_pdf(content)
            
            # Planilhas
            elif extension in ['csv', 'txt']:
                return self._extract_csv(content)
            elif extension in ['xlsx', 'xls']:
                return self._extract_excel(content)
            
            # Imagens (OCR)
            elif extension in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'tiff']:
                return self._extract_image_ocr(content)
            
            # Word
            elif extension == 'docx':
                return self._extract_docx(content)
            elif extension == 'doc':
                return self._extract_doc_legacy(content)
            
            # PowerPoint
            elif extension == 'pptx':
                return self._extract_pptx(content)
            elif extension == 'ppt':
                return self._extract_ppt_legacy(content)
            
            else:
                raise ValueError(f"Formato não suportado: {extension}")
                
        except Exception as e:
            logger.error(f"Erro ao extrair texto de {filename}: {e}")
            raise
    
    # ========================================================================
    # PDF
    # ========================================================================
    
    def _extract_pdf(self, content: bytes) -> str:
        """Extrai texto de PDF usando pdfplumber (melhor para tabelas)"""
        text_parts = []
        
        try:
            # Tentar com pdfplumber primeiro (melhor qualidade)
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            
            if text_parts:
                return '\n\n'.join(text_parts)
        except Exception as e:
            logger.warning(f"pdfplumber falhou, tentando PyPDF2: {e}")
        
        # Fallback para PyPDF2
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return '\n\n'.join(text_parts)
        except Exception as e:
            logger.error(f"PyPDF2 também falhou: {e}")
            raise ValueError("Não foi possível extrair texto do PDF")
    
    # ========================================================================
    # CSV / TXT
    # ========================================================================
    
    def _extract_csv(self, content: bytes) -> str:
        """Extrai texto de CSV"""
        try:
            # Tentar UTF-8
            text = content.decode('utf-8')
        except UnicodeDecodeError:
            # Fallback para latin1
            text = content.decode('latin1')
        
        # Converter CSV para texto legível
        try:
            df = pd.read_csv(io.StringIO(text))
            # Converter DataFrame para texto formatado
            return df.to_string(index=False)
        except:
            # Se falhar, retornar texto bruto
            return text
    
    # ========================================================================
    # Excel
    # ========================================================================
    
    def _extract_excel(self, content: bytes) -> str:
        """Extrai texto de Excel"""
        try:
            # Ler todas as abas
            excel_file = pd.ExcelFile(io.BytesIO(content))
            
            if len(excel_file.sheet_names) == 1:
                # Uma aba só
                df = pd.read_excel(io.BytesIO(content))
                return df.to_string(index=False)
            else:
                # Múltiplas abas
                text_parts = []
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(io.BytesIO(content), sheet_name=sheet_name)
                    text_parts.append(f"=== {sheet_name} ===\n{df.to_string(index=False)}")
                return '\n\n'.join(text_parts)
                
        except Exception as e:
            logger.error(f"Erro ao processar Excel: {e}")
            raise
    
    # ========================================================================
    # IMAGENS - OCR
    # ========================================================================
    
    def _extract_image_ocr(self, content: bytes) -> str:
        """
        Extrai texto de imagem usando Tesseract OCR
        
        Suporta: JPG, PNG, GIF, BMP, WEBP, TIFF
        """
        try:
            # Abrir imagem
            image = Image.open(io.BytesIO(content))
            
            # Converter para RGB se necessário
            if image.mode not in ('RGB', 'L'):
                image = image.convert('RGB')
            
            # Rotacionar se necessário (alguns smartphones)
            try:
                from PIL import ImageOps
                image = ImageOps.exif_transpose(image)
            except:
                pass
            
            # Extrair texto com Tesseract
            # lang='por+eng' = Português + Inglês
            # psm 6 = Assume bloco uniforme de texto
            text = pytesseract.image_to_string(
                image,
                lang='por+eng',
                config='--psm 6 --oem 3'
            )
            
            # Limpar texto
            text = text.strip()
            
            if not text or len(text) < 5:
                raise ValueError("Nenhum texto significativo detectado na imagem")
            
            logger.info(f"✅ OCR extraiu {len(text)} caracteres da imagem")
            return text
            
        except Exception as e:
            logger.error(f"Erro no OCR: {e}")
            raise ValueError(f"Não foi possível extrair texto da imagem: {e}")
    
    # ========================================================================
    # WORD - DOCX (novo formato)
    # ========================================================================
    
    def _extract_docx(self, content: bytes) -> str:
        """
        Extrai texto de arquivo Word (.docx)
        
        Formato: Office 2007+
        """
        try:
            doc = Document(io.BytesIO(content))
            
            text_parts = []
            
            # Extrair parágrafos
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extrair tabelas
            for table in doc.tables:
                table_text = self._extract_table_from_word(table)
                if table_text:
                    text_parts.append(table_text)
            
            text = '\n\n'.join(text_parts)
            
            if not text:
                raise ValueError("Documento Word vazio")
            
            logger.info(f"✅ Extraído {len(text)} caracteres do DOCX")
            return text
            
        except Exception as e:
            logger.error(f"Erro ao processar DOCX: {e}")
            raise ValueError(f"Não foi possível extrair texto do Word: {e}")
    
    def _extract_table_from_word(self, table) -> str:
        """Extrai texto de tabela do Word"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(' | '.join(cells))
        return '\n'.join(rows)
    
    # ========================================================================
    # WORD - DOC (formato legado)
    # ========================================================================
    
    def _extract_doc_legacy(self, content: bytes) -> str:
        """
        Extrai texto de arquivo Word legado (.doc)
        
        Formato: Office 97-2003
        Usa LibreOffice para converter para texto
        """
        try:
            # Criar arquivo temporário
            with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_input:
                tmp_input.write(content)
                tmp_input_path = tmp_input.name
            
            # Criar arquivo de saída temporário
            tmp_output_path = tmp_input_path.replace('.doc', '.txt')
            
            try:
                # Converter com LibreOffice
                subprocess.run([
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'txt:Text',
                    '--outdir', os.path.dirname(tmp_input_path),
                    tmp_input_path
                ], check=True, timeout=30, capture_output=True)
                
                # Ler resultado
                with open(tmp_output_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                
                if not text or len(text) < 5:
                    raise ValueError("Documento DOC vazio ou conversão falhou")
                
                logger.info(f"✅ Extraído {len(text)} caracteres do DOC legado")
                return text
                
            finally:
                # Limpar arquivos temporários
                try:
                    os.unlink(tmp_input_path)
                    if os.path.exists(tmp_output_path):
                        os.unlink(tmp_output_path)
                except:
                    pass
                    
        except subprocess.TimeoutExpired:
            raise ValueError("Timeout ao converter DOC legado")
        except Exception as e:
            logger.error(f"Erro ao processar DOC legado: {e}")
            raise ValueError(f"Não foi possível extrair texto do Word legado: {e}")
    
    # ========================================================================
    # POWERPOINT - PPTX (novo formato)
    # ========================================================================
    
    def _extract_pptx(self, content: bytes) -> str:
        """
        Extrai texto de arquivo PowerPoint (.pptx)
        
        Formato: Office 2007+
        """
        try:
            prs = Presentation(io.BytesIO(content))
            
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"=== Slide {slide_num} ==="
                
                # Extrair texto de todas as shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"\n{shape.text}"
                    
                    # Extrair texto de tabelas
                    if hasattr(shape, "table"):
                        table_text = self._extract_table_from_ppt(shape.table)
                        if table_text:
                            slide_text += f"\n{table_text}"
                
                if len(slide_text) > len(f"=== Slide {slide_num} ==="):
                    text_parts.append(slide_text)
            
            text = '\n\n'.join(text_parts)
            
            if not text:
                raise ValueError("PowerPoint vazio")
            
            logger.info(f"✅ Extraído {len(text)} caracteres de {len(prs.slides)} slides")
            return text
            
        except Exception as e:
            logger.error(f"Erro ao processar PPTX: {e}")
            raise ValueError(f"Não foi possível extrair texto do PowerPoint: {e}")
    
    def _extract_table_from_ppt(self, table) -> str:
        """Extrai texto de tabela do PowerPoint"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(' | '.join(cells))
        return '\n'.join(rows)
    
    # ========================================================================
    # POWERPOINT - PPT (formato legado)
    # ========================================================================
    
    def _extract_ppt_legacy(self, content: bytes) -> str:
        """
        Extrai texto de arquivo PowerPoint legado (.ppt)
        
        Formato: Office 97-2003
        Usa LibreOffice para converter para texto
        """
        try:
            # Criar arquivo temporário
            with tempfile.NamedTemporaryFile(suffix='.ppt', delete=False) as tmp_input:
                tmp_input.write(content)
                tmp_input_path = tmp_input.name
            
            # Criar arquivo de saída temporário
            tmp_output_path = tmp_input_path.replace('.ppt', '.txt')
            
            try:
                # Converter com LibreOffice
                subprocess.run([
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'txt:Text',
                    '--outdir', os.path.dirname(tmp_input_path),
                    tmp_input_path
                ], check=True, timeout=30, capture_output=True)
                
                # Ler resultado
                with open(tmp_output_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                
                if not text or len(text) < 5:
                    raise ValueError("PowerPoint PPT vazio ou conversão falhou")
                
                logger.info(f"✅ Extraído {len(text)} caracteres do PPT legado")
                return text
                
            finally:
                # Limpar arquivos temporários
                try:
                    os.unlink(tmp_input_path)
                    if os.path.exists(tmp_output_path):
                        os.unlink(tmp_output_path)
                except:
                    pass
                    
        except subprocess.TimeoutExpired:
            raise ValueError("Timeout ao converter PPT legado")
        except Exception as e:
            logger.error(f"Erro ao processar PPT legado: {e}")
            raise ValueError(f"Não foi possível extrair texto do PowerPoint legado: {e}")